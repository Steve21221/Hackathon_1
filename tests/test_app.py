import json
import os
import tempfile
import unittest
from io import BytesIO
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import Mock, patch

from app import (
    app,
    build_feedback_prompt,
    call_model,
    extract_generated_prompt,
    fetch_web_source,
    feedback_timing_profile,
    load_mentor_prompt,
    normalize_public_web_url,
    record_feedback_performance,
    render_markdown_html,
    split_review_text,
)
from docx import Document
from PIL import Image
from pptx import Presentation


class PromptlyTestCase(unittest.TestCase):
    def setUp(self):
        os.environ.pop("MODEL_PROVIDER", None)
        os.environ.pop("OLLAMA_BASE_URL", None)
        os.environ.pop("OLLAMA_MODEL", None)
        os.environ.pop("OPENAI_API_KEY", None)
        os.environ.pop("OPENAI_MODEL", None)
        os.environ.pop("ANTHROPIC_API_KEY", None)
        os.environ.pop("CLAUDE_MODEL", None)
        app.config.update(TESTING=True)
        self.original_output_directory = app.config["OUTPUT_DIR"]
        self.original_mentor_library_directory = app.config["MENTOR_LIBRARY_DIR"]
        self.original_settings_path = app.config.get("SETTINGS_PATH")
        self.original_performance_path = app.config.get("PERFORMANCE_PATH")
        self.output_temp = tempfile.TemporaryDirectory()
        app.config["OUTPUT_DIR"] = Path(self.output_temp.name) / "outputs"
        app.config["MENTOR_LIBRARY_DIR"] = Path(self.output_temp.name) / "mentor_files"
        app.config["SETTINGS_PATH"] = Path(self.output_temp.name) / "user_settings.json"
        app.config["PERFORMANCE_PATH"] = Path(self.output_temp.name) / "model_performance.json"
        self.client = app.test_client()

    def tearDown(self):
        app.config["OUTPUT_DIR"] = self.original_output_directory
        app.config["MENTOR_LIBRARY_DIR"] = self.original_mentor_library_directory
        if self.original_settings_path is not None:
            app.config["SETTINGS_PATH"] = self.original_settings_path
        if self.original_performance_path is not None:
            app.config["PERFORMANCE_PATH"] = self.original_performance_path
        self.output_temp.cleanup()

    def test_home_page_has_three_feedback_categories(self):
        response = self.client.get("/")
        self.assertEqual(response.status_code, 200)
        self.assertIn(b"Papers &amp; proposals", response.data)
        self.assertIn(b"Research ideas", response.data)
        self.assertIn(b"Talks &amp; slides", response.data)
        self.assertIn(b'data-choose-prompt', response.data)
        self.assertIn(b'data-upload-panel="research-ideas"', response.data)
        self.assertEqual(response.data.count(b"data-feedback-file-input"), 3)
        self.assertEqual(response.data.count(b"data-feedback-file-selection"), 3)
        self.assertEqual(response.data.count(b"data-remove-feedback-file"), 3)
        self.assertEqual(response.data.count(b"data-web-research-toggle"), 3)
        self.assertEqual(response.data.count(b"data-web-source-urls"), 3)

    def test_upload_script_supports_removing_pending_files(self):
        script = (Path("static") / "library_uploads.js").read_text(encoding="utf-8")

        self.assertIn("removeSelectedReferenceFile", script)
        self.assertIn("_promptlyPendingFiles", script)
        self.assertIn("body.delete(input.name)", script)
        self.assertIn("data-remove-feedback-file", script)
        self.assertIn("input.value = '';", script)
        self.assertIn("data-restore-defaults-form", script)
        self.assertIn("downloaded models stay unchanged", script)
        self.assertIn("Calibrating\\u2026", script)
        self.assertIn("Estimated response time:", script)
        self.assertIn("feedback_timing", Path("app.py").read_text(encoding="utf-8"))
        self.assertIn("data-web-research-toggle", script)
        self.assertIn("urls.required = enabled", script)

    def test_feedback_timing_calibrates_separately_for_each_model(self):
        os.environ["MODEL_PROVIDER"] = "ollama"
        os.environ["OLLAMA_MODEL"] = "qwen3.5:9b"

        self.assertEqual(feedback_timing_profile()["state"], "calibrating")
        for index in range(10):
            record_feedback_performance(
                provider="ollama",
                file_bytes=1_000 + index,
                extracted_chars=2_000 + index,
                elapsed_seconds=60 + index,
            )

        calibrated = feedback_timing_profile()
        self.assertEqual(calibrated["state"], "calibrated")
        self.assertEqual(calibrated["model"], "qwen3.5:9b")
        self.assertEqual(calibrated["sample_count"], 8)
        self.assertGreater(calibrated["average_seconds"], 60)
        self.assertTrue(Path(app.config["PERFORMANCE_PATH"]).is_file())

        os.environ["OLLAMA_MODEL"] = "phi4-mini"
        phi_profile = feedback_timing_profile()
        self.assertEqual(phi_profile["state"], "calibrating")
        self.assertEqual(phi_profile["model"], "phi4-mini")

    @patch("app.call_model", return_value="Measured mentor feedback")
    def test_successful_feedback_request_completes_calibration(self, _mock_model):
        os.environ["MODEL_PROVIDER"] = "ollama"
        os.environ["OLLAMA_MODEL"] = "qwen3.5:4b"
        home = self.client.get("/")
        self.assertIn(b'"state": "calibrating"', home.data)

        response = self.client.post(
            "/feedback",
            data={
                "content_type": "research-ideas",
                "mentor_id": "dr-nanshu-lu",
                "file": (BytesIO(b"A measurable research idea."), "idea.txt"),
            },
            content_type="multipart/form-data",
            headers={"X-Requested-With": "XMLHttpRequest", "Accept": "application/json"},
        )

        self.assertEqual(response.status_code, 200)
        timing = response.get_json()["feedback_timing"]
        self.assertEqual(timing["state"], "calibrated")
        self.assertEqual(timing["model"], "qwen3.5:4b")
        self.assertEqual(timing["sample_count"], 1)

    def test_review_style_workspace_is_separate_from_home_page(self):
        response = self.client.get("/")
        self.assertEqual(response.status_code, 200)
        self.assertIn(b'href="/prompt-library"', response.data)
        self.assertIn(b"Modify a review style", response.data)
        self.assertLess(
            response.data.index(b"Modify a review style"),
            response.data.index(b"Model settings"),
        )
        self.assertIn(b"What kind of feedback do you need?", response.data)
        self.assertNotIn(b'name="research_files"', response.data)
        self.assertNotIn(b"Generate reusable prompts", response.data)
        self.assertIn(b"library_uploads.js", response.data)
        self.assertNotIn(b"Mentor data", response.data)

        library_response = self.client.get("/prompt-library")
        self.assertEqual(library_response.status_code, 200)
        self.assertIn(b"Modify a review style", library_response.data)
        self.assertIn(b'name="research_files"', library_response.data)
        self.assertIn(b'name="slide_files"', library_response.data)
        self.assertIn(b'name="paper_files"', library_response.data)
        self.assertIn(b"Generate reusable prompts", library_response.data)
        self.assertIn(b"Feedback workspace", library_response.data)
        self.assertLess(
            library_response.data.index(b"Papers / proposals"),
            library_response.data.index(b"Research ideas / meeting minutes"),
        )
        self.assertLess(
            library_response.data.index(b"Research ideas / meeting minutes"),
            library_response.data.index(b"Talks / presentations / slides"),
        )
        self.assertIn(b"reference-choice-paper-proposal-pi", library_response.data)
        self.assertIn(b"reference-choice-meeting-research-pi", library_response.data)
        self.assertIn(b"reference-choice-slides-talk-pi", library_response.data)
        self.assertLess(
            library_response.data.index(b"Feedback workspace"),
            library_response.data.index(b"Model settings"),
        )
        self.assertNotIn(b"What kind of feedback do you need?", library_response.data)

        self.assertEqual(self.client.get("/mentor-data").status_code, 404)

    def test_home_page_shows_static_starting_mentor_card(self):
        response = self.client.get("/?type=research-ideas")
        self.assertEqual(response.status_code, 200)
        self.assertIn(b"Dr. Nanshu Lu", response.data)
        self.assertIn(b"Available mentor", response.data)
        self.assertIn(b"Built-in starting profile", response.data)
        self.assertNotIn(b"No mentor libraries yet", response.data)
        self.assertIn(b"Modify a review style", response.data)

    def test_prompt_generation_returns_review_style_page_with_ready_prompts(self):
        response = self.client.post(
            "/generate-prompts",
            data={
                "research_files": (
                    BytesIO(
                        b"Professor asks us to clarify the hypothesis, compare mechanisms, "
                        b"and define the next experiment."
                    ),
                    "meeting-notes.txt",
                )
            },
            content_type="multipart/form-data",
        )

        self.assertEqual(response.status_code, 200)
        self.assertIn(b"Promptly &mdash; PI-style prompt library", response.data)
        self.assertIn(b"PI-style prompts ready", response.data)
        self.assertIn(b"Download TXT", response.data)
        self.assertNotIn(b"What kind of feedback do you need?", response.data)

    def test_pi_style_library_generates_and_downloads_prompt(self):
        response = self.client.post(
            "/generate-prompts",
            data={
                "research_files": (
                    BytesIO(
                        b"Professor asks us to clarify the hypothesis, compare mechanisms, "
                        b"and define the next experiment."
                    ),
                    "meeting-notes.txt",
                )
            },
            content_type="multipart/form-data",
        )

        self.assertEqual(response.status_code, 200)
        self.assertIn(b"PI-style prompts ready", response.data)
        self.assertIn(b"Research Ideas / Meeting Minutes", response.data)
        self.assertNotIn(b"Talks / Presentations / Slides</h3>", response.data)
        run_id = response.headers["X-Prompt-Run-Id"]
        run_directory = Path(app.config["OUTPUT_DIR"]) / run_id
        self.assertTrue((run_directory / "meeting_research_pi_prompt.txt").is_file())
        self.assertTrue((run_directory / "all_pi_style_prompts.txt").is_file())
        run_metadata = json.loads(
            (run_directory / ".mentor-run.json").read_text(encoding="utf-8")
        )
        self.assertEqual(run_metadata["mentor_slug"], "pi-style-library")

        download = self.client.get(f"/download/{run_id}/meeting_research_pi_prompt")
        self.assertEqual(download.status_code, 200)
        self.assertIn(b"Generated PI-style prompt", download.data)
        download.close()

        combined = (run_directory / "all_pi_style_prompts.txt").read_text(encoding="utf-8")
        self.assertIn("SOURCE FILES: meeting-notes.txt", combined)
        self.assertNotIn("No uploaded references; using default mode prompt.", combined)
        self.assertNotIn("MODE: slides_talk_pi", combined)
        self.assertNotIn("MODE: paper_proposal_pi", combined)

    @patch("app.call_prompt_model")
    def test_pi_style_library_uses_configured_model_for_style_distillation(self, mock_prompt_model):
        os.environ["MODEL_PROVIDER"] = "ollama"
        mock_prompt_model.side_effect = [
            "- clarify the framing\n- compare mechanisms\n- convert discussion into actions",
            "Review future research by reframing the opportunity, comparing mechanisms, and ending with decisive actions.",
        ]

        response = self.client.post(
            "/generate-prompts",
            data={"research_files": (BytesIO(b"A detailed professor review example."), "review.txt")},
            content_type="multipart/form-data",
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(mock_prompt_model.call_count, 2)
        run_id = response.headers["X-Prompt-Run-Id"]
        generated = (Path(app.config["OUTPUT_DIR"]) / run_id / "meeting_research_pi_prompt.txt").read_text(
            encoding="utf-8"
        )
        self.assertIn("Review future research by reframing the opportunity", generated)

    def test_pi_style_library_rejects_unsupported_reference_file(self):
        response = self.client.post(
            "/generate-prompts",
            data={"research_files": (BytesIO(b"binary"), "unsafe.exe")},
            content_type="multipart/form-data",
        )
        self.assertEqual(response.status_code, 400)
        self.assertIn(b"not supported", response.data)

    def test_existing_prompt_libraries_use_mentor_card_selector_on_review_style_page(self):
        self.client.post(
            "/generate-prompts",
            data={
                "prompt_mentor_name": "Card Mentor",
                "research_files": (BytesIO(b"Ask why the mechanism is testable."), "card.txt"),
            },
            content_type="multipart/form-data",
        )

        response = self.client.get("/prompt-library?prompt_mentor=card-mentor")

        self.assertEqual(response.status_code, 200)
        self.assertIn(b"Choose your mentor", response.data)
        self.assertIn(b"Dr. Nanshu Lu", response.data)
        self.assertIn(b'value="dr-nanshu-lu"', response.data)
        self.assertIn(b"Card Mentor", response.data)
        self.assertIn(b"PI-style library", response.data)
        self.assertIn(b'class="mentor-card selected"', response.data)
        self.assertIn(b'name="selected_prompt_mentor" value="card-mentor"', response.data)
        self.assertNotIn(b"Select an existing library", response.data)
        self.assertNotIn(b'<select name="selected_prompt_mentor"', response.data)
        self.assertLess(
            response.data.index(b"Choose your mentor"),
            response.data.index(b"Upload examples by review category"),
        )

    def test_delete_mentor_uses_a_separate_form_from_new_mentor_creation(self):
        response = self.client.get("/prompt-library")

        self.assertEqual(response.status_code, 200)
        self.assertIn(b'id="delete-mentor-form"', response.data)
        self.assertIn(b'form="delete-mentor-form"', response.data)
        self.assertIn(b'data-delete-mentor-form', response.data)
        self.assertIn(b'name="selected_prompt_mentor"', response.data)
        self.assertNotIn(b'formaction="/delete-mentor"', response.data)
        self.assertIn(b'action="/restore-mentor-defaults"', response.data)
        self.assertIn(b'data-restore-defaults-form', response.data)
        self.assertIn(b"Restore mentor defaults", response.data)
        self.assertIn(b"downloaded models stay unchanged", response.data)

    def test_builtin_mentor_library_is_shared_with_feedback_workspace(self):
        selected = self.client.get("/prompt-library?prompt_mentor=dr-nanshu-lu")
        self.assertEqual(selected.status_code, 200)
        self.assertIn(b"Dr. Nanshu Lu", selected.data)
        self.assertIn(b"Built-in mentor", selected.data)

        library_api = self.client.get("/api/library/dr-nanshu-lu")
        self.assertEqual(library_api.status_code, 200)
        self.assertEqual(library_api.get_json()["name"], "Dr. Nanshu Lu")

        created = self.client.post(
            "/generate-prompts",
            data={
                "selected_prompt_mentor": "dr-nanshu-lu",
                "slide_files": (
                    BytesIO(b"Make the central takeaway explicit and remove duplicate panels."),
                    "local-slides.txt",
                ),
            },
            content_type="multipart/form-data",
        )
        self.assertEqual(created.status_code, 200)
        local_prompt = (
            Path(app.config["MENTOR_LIBRARY_DIR"])
            / "dr-nanshu-lu"
            / "slides_talk_pi"
            / "prompt.txt"
        )
        self.assertTrue(local_prompt.is_file())

        self.assertEqual(
            load_mentor_prompt("dr-nanshu-lu", "talks-slides"),
            extract_generated_prompt(local_prompt.read_text(encoding="utf-8")),
        )
        self.assertIn(
            "falsifiable",
            load_mentor_prompt("dr-nanshu-lu", "research-ideas").lower(),
        )

        home = self.client.get("/")
        self.assertEqual(home.status_code, 200)
        self.assertIn(b"Built-in profile with local updates", home.data)

    def test_named_prompt_library_persists_files_and_stable_prompts(self):
        response = self.client.post(
            "/generate-prompts",
            data={
                "prompt_mentor_name": "Dr. Custom Mentor",
                "research_files": (
                    BytesIO(b"Clarify the hypothesis and define the next experiment."),
                    "meeting-notes.txt",
                ),
            },
            content_type="multipart/form-data",
        )

        self.assertEqual(response.status_code, 200)
        library = Path(app.config["MENTOR_LIBRARY_DIR"]) / "dr-custom-mentor"
        self.assertTrue((library / "mentor.json").is_file())
        self.assertTrue(
            (library / "meeting_research_pi" / "raw" / "meeting-notes.txt").is_file()
        )
        self.assertTrue((library / "meeting_research_pi" / "prompt.txt").is_file())
        self.assertTrue((library / "all_pi_style_prompts.txt").is_file())
        self.assertIn(b"Dr. Custom Mentor", response.data)
        self.assertIn(b"meeting-notes.txt", response.data)

    def test_generated_prompts_explain_stable_txt_files_are_updated(self):
        response = self.client.post(
            "/generate-prompts",
            data={
                "prompt_mentor_name": "Stable Mentor",
                "research_files": (BytesIO(b"Clarify the central hypothesis."), "stable.txt"),
            },
            content_type="multipart/form-data",
        )

        self.assertEqual(response.status_code, 200)
        self.assertIn(b"Saved local TXT files to", response.data)
        self.assertIn(b"Regenerating this same PI library and mode updates the existing TXT file", response.data)

        stable_prompt = Path(app.config["MENTOR_LIBRARY_DIR"]) / "stable-mentor" / "meeting_research_pi" / "prompt.txt"
        self.assertTrue(stable_prompt.is_file())
        first_path = stable_prompt.resolve()
        first_content = stable_prompt.read_text(encoding="utf-8")

        self.client.post(
            "/generate-prompts",
            data={
                "selected_prompt_mentor": "stable-mentor",
                "research_files": (BytesIO(b"Define the next decisive control."), "stable.txt"),
            },
            content_type="multipart/form-data",
        )

        self.assertEqual(stable_prompt.resolve(), first_path)
        self.assertNotEqual(stable_prompt.read_text(encoding="utf-8"), first_content)

    def test_stored_reference_file_can_be_deleted_from_web(self):
        self.client.post(
            "/generate-prompts",
            data={
                "prompt_mentor_name": "File Delete Mentor",
                "paper_files": (BytesIO(b"Tighten the proposal argument."), "delete-me.txt"),
            },
            content_type="multipart/form-data",
        )
        stored_file = Path(app.config["MENTOR_LIBRARY_DIR"]) / "file-delete-mentor" / "paper_proposal_pi" / "raw" / "delete-me.txt"
        self.assertTrue(stored_file.is_file())

        selected = self.client.get("/prompt-library?prompt_mentor=file-delete-mentor")
        self.assertEqual(selected.status_code, 200)
        self.assertIn(b"delete-me.txt", selected.data)
        self.assertIn(b'name="delete_reference_file" value="paper_proposal_pi|delete-me.txt"', selected.data)

        deleted = self.client.post(
            "/delete-reference-file",
            data={
                "selected_prompt_mentor": "file-delete-mentor",
                "delete_reference_file": "paper_proposal_pi|delete-me.txt",
            },
            follow_redirects=True,
        )

        self.assertEqual(deleted.status_code, 200)
        self.assertFalse(stored_file.exists())
        self.assertNotIn(b"delete-me.txt", deleted.data)
        self.assertIn(b"File Delete Mentor", deleted.data)

    def test_delete_stored_reference_file_rejects_path_escape(self):
        response = self.client.post(
            "/delete-reference-file",
            data={
                "selected_prompt_mentor": "../outside",
                "delete_reference_file": "paper_proposal_pi|../secret.txt",
            },
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn(b"Please select an existing library file to delete", response.data)

    def test_existing_prompt_library_can_regenerate_without_new_uploads(self):
        self.client.post(
            "/generate-prompts",
            data={
                "prompt_mentor_name": "Existing Mentor",
                "research_files": (BytesIO(b"Ask for missing controls."), "first.txt"),
            },
            content_type="multipart/form-data",
        )

        response = self.client.post(
            "/generate-prompts",
            data={"selected_prompt_mentor": "existing-mentor"},
            content_type="multipart/form-data",
        )

        self.assertEqual(response.status_code, 200)
        self.assertIn(b"PI-style prompts ready", response.data)
        self.assertIn(b"first.txt", response.data)

    def test_prompt_library_selection_and_delete_flow(self):
        created = self.client.post(
            "/generate-prompts",
            data={
                "prompt_mentor_name": "Delete Mentor",
                "paper_files": (BytesIO(b"Tighten the argument."), "review.txt"),
            },
            content_type="multipart/form-data",
        )
        self.assertEqual(created.status_code, 200)
        run_id = created.headers["X-Prompt-Run-Id"]
        run_directory = Path(app.config["OUTPUT_DIR"]) / run_id
        library = Path(app.config["MENTOR_LIBRARY_DIR"]) / "delete-mentor"
        unrelated_run = Path(app.config["OUTPUT_DIR"]) / "pi_style_prompts_other-mentor_keep"
        unrelated_run.mkdir()
        (unrelated_run / ".mentor-run.json").write_text(
            json.dumps({"mentor_slug": "other-mentor"}),
            encoding="utf-8",
        )
        self.assertTrue(run_directory.is_dir())

        selected = self.client.get("/prompt-library?prompt_mentor=delete-mentor")
        self.assertEqual(selected.status_code, 200)
        self.assertIn(b"Delete Mentor", selected.data)
        self.assertIn(b"Delete mentor", selected.data)
        self.assertIn(b"review.txt", selected.data)

        deleted = self.client.post(
            "/delete-mentor",
            data={"selected_prompt_mentor": "delete-mentor"},
            follow_redirects=True,
        )
        self.assertEqual(deleted.status_code, 200)
        self.assertFalse(library.exists())
        self.assertFalse(run_directory.exists())
        self.assertTrue(unrelated_run.is_dir())
        self.assertNotIn(b"review.txt", deleted.data)

    def test_builtin_mentor_can_be_deleted_and_restored(self):
        created = self.client.post(
            "/generate-prompts",
            data={
                "selected_prompt_mentor": "dr-nanshu-lu",
                "paper_files": (
                    BytesIO(b"State the hypothesis and compare against a strong baseline."),
                    "local-paper.txt",
                ),
            },
            content_type="multipart/form-data",
        )
        self.assertEqual(created.status_code, 200)
        run_directory = Path(app.config["OUTPUT_DIR"]) / created.headers["X-Prompt-Run-Id"]
        library = Path(app.config["MENTOR_LIBRARY_DIR"]) / "dr-nanshu-lu"
        self.assertTrue(run_directory.is_dir())
        self.assertTrue(library.is_dir())

        response = self.client.post(
            "/delete-mentor",
            data={"selected_prompt_mentor": "dr-nanshu-lu"},
            follow_redirects=True,
        )

        self.assertEqual(response.status_code, 200)
        self.assertNotIn(b"Dr. Nanshu Lu", response.data)
        self.assertFalse(library.exists())
        self.assertFalse(run_directory.exists())
        self.assertTrue((Path("Mentor_Data") / "dr-nanshu-lu").is_dir())
        removed_path = Path(app.config["MENTOR_LIBRARY_DIR"]) / ".removed-mentors.json"
        self.assertEqual(json.loads(removed_path.read_text(encoding="utf-8")), ["dr-nanshu-lu"])
        self.assertNotIn(b"Dr. Nanshu Lu", self.client.get("/").data)
        self.assertEqual(self.client.get("/api/library/dr-nanshu-lu").status_code, 404)
        with self.assertRaisesRegex(ValueError, "available mentor"):
            load_mentor_prompt("dr-nanshu-lu", "papers-proposals")

        restored = self.client.post(
            "/generate-prompts",
            data={
                "prompt_mentor_name": "Dr. Nanshu Lu",
                "research_files": (
                    BytesIO(b"Ask what observation would falsify the central idea."),
                    "restored.txt",
                ),
            },
            content_type="multipart/form-data",
        )
        self.assertEqual(restored.status_code, 200)
        self.assertIn(b"Dr. Nanshu Lu", restored.data)
        self.assertFalse(removed_path.exists())
        self.assertIn(b"Dr. Nanshu Lu", self.client.get("/").data)

    def test_restore_mentor_defaults_clears_only_local_mentor_state(self):
        mentor_root = Path(app.config["MENTOR_LIBRARY_DIR"])
        output_root = Path(app.config["OUTPUT_DIR"])
        settings_path = Path(app.config["SETTINGS_PATH"])
        local_library = mentor_root / "custom-mentor"
        local_override = mentor_root / "dr-nanshu-lu" / "paper_proposal_pi"
        local_library.mkdir(parents=True)
        local_override.mkdir(parents=True)
        output_run = output_root / "pi_style_prompts_custom-mentor_test"
        output_run.mkdir(parents=True)
        (local_library / "mentor.json").write_text(
            json.dumps({"slug": "custom-mentor", "name": "Custom Mentor"}),
            encoding="utf-8",
        )
        (local_override / "prompt.txt").write_text("Local override", encoding="utf-8")
        (mentor_root / ".removed-mentors.json").write_text(
            json.dumps(["dr-nanshu-lu"]),
            encoding="utf-8",
        )
        (output_run / "all_pi_style_prompts.txt").write_text(
            "Local run copy",
            encoding="utf-8",
        )
        saved_settings = json.dumps(
            {
                "MODEL_PROVIDER": "ollama",
                "OLLAMA_MODEL": "qwen3.5:4b",
                "OPENAI_API_KEY": "preserve-me",
            },
            indent=2,
        )
        settings_path.write_text(saved_settings, encoding="utf-8")

        response = self.client.post(
            "/restore-mentor-defaults",
            follow_redirects=True,
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(list(mentor_root.iterdir()), [])
        self.assertEqual(list(output_root.iterdir()), [])
        self.assertEqual(settings_path.read_text(encoding="utf-8"), saved_settings)
        self.assertIn(b"Mentor defaults restored", response.data)
        self.assertIn(b"Dr. Nanshu Lu", response.data)
        self.assertNotIn(b"Custom Mentor", response.data)
        self.assertTrue(
            (Path("Mentor_Data") / "dr-nanshu-lu" / "paper_proposal_pi.txt").is_file()
        )

    def test_delete_prompt_library_rejects_unsafe_selection(self):
        response = self.client.post(
            "/delete-mentor",
            data={"selected_prompt_mentor": "../outside"},
        )
        self.assertEqual(response.status_code, 400)
        self.assertIn(b"existing mentor", response.data)

    def test_clicking_type_shows_one_upload_form(self):
        response = self.client.get("/?type=research-ideas")
        self.assertEqual(response.status_code, 200)
        self.assertIn(b'data-upload-panel="research-ideas"', response.data)
        self.assertIn(b"Choose a research idea", response.data)
        self.assertIn(b'data-type-choice', response.data)
        self.assertIn(b'data-feedback-mentor-card', response.data)

    def test_research_idea_upload_returns_demo_feedback(self):
        response = self.client.post(
            "/feedback",
            data={
                "content_type": "research-ideas",
                "focus": "key decisions",
                "file": (BytesIO(b"Speaker one: We approved the project."), "meeting.txt"),
            },
            content_type="multipart/form-data",
        )
        self.assertEqual(response.status_code, 200)
        self.assertIn(b"uploaded and read successfully", response.data)
        self.assertIn(b"markdown-body", response.data)

    def test_ajax_feedback_returns_json_without_page_reload_payload(self):
        response = self.client.post(
            "/feedback",
            data={
                "content_type": "research-ideas",
                "mentor_id": "dr-nanshu-lu",
                "file": (BytesIO(b"An early research idea."), "idea.txt"),
            },
            content_type="multipart/form-data",
            headers={"X-Requested-With": "XMLHttpRequest", "Accept": "application/json"},
        )
        self.assertEqual(response.status_code, 200)
        payload = response.get_json()
        self.assertIn("output_html", payload)
        self.assertIn("<h2>", payload["output_html"])
        self.assertIn("idea.txt", payload["filename"])

    @patch("app.call_model", return_value="Feedback with [Source 1](https://example.com/research)")
    @patch("app.load_web_sources")
    def test_feedback_with_web_sources_adds_citation_context(self, mock_load_sources, mock_model):
        mock_load_sources.return_value = [
            {
                "title": "Research reference",
                "url": "https://example.com/research",
                "text": "A relevant public research summary.",
            }
        ]

        response = self.client.post(
            "/feedback",
            data={
                "content_type": "research-ideas",
                "mentor_id": "dr-nanshu-lu",
                "use_web_sources": "1",
                "web_source_urls": "https://example.com/research",
                "file": (BytesIO(b"An early research idea."), "idea.txt"),
            },
            content_type="multipart/form-data",
            headers={"X-Requested-With": "XMLHttpRequest", "Accept": "application/json"},
        )

        self.assertEqual(response.status_code, 200)
        mock_load_sources.assert_called_once_with("https://example.com/research")
        prompt = mock_model.call_args.args[0]
        self.assertIn("untrusted reference content", prompt)
        self.assertIn("[Source 1] Research reference", prompt)
        self.assertIn("[Source 1](https://example.com/research)", prompt)
        self.assertIn("https://example.com/research", prompt)
        self.assertEqual(
            response.get_json()["web_sources"],
            [{"title": "Research reference", "url": "https://example.com/research"}],
        )

    @patch("app.load_web_sources")
    def test_feedback_does_not_fetch_urls_without_opt_in(self, mock_load_sources):
        response = self.client.post(
            "/feedback",
            data={
                "content_type": "research-ideas",
                "mentor_id": "dr-nanshu-lu",
                "web_source_urls": "https://example.com/ignored",
                "file": (BytesIO(b"An early research idea."), "idea.txt"),
            },
            content_type="multipart/form-data",
        )

        self.assertEqual(response.status_code, 200)
        mock_load_sources.assert_not_called()

    def test_feedback_requires_url_when_web_sources_are_enabled(self):
        response = self.client.post(
            "/feedback",
            data={
                "content_type": "research-ideas",
                "mentor_id": "dr-nanshu-lu",
                "use_web_sources": "1",
                "file": (BytesIO(b"An early research idea."), "idea.txt"),
            },
            content_type="multipart/form-data",
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn(b"Add at least one public URL", response.data)

    def test_web_source_url_validation_blocks_private_targets_and_credentials(self):
        blocked_urls = (
            "http://127.0.0.1/private",
            "http://localhost/private",
            "http://192.168.1.10/private",
            "https://user:password@example.com/private",
            "https://example.com:11434/private",
        )
        for source_url in blocked_urls:
            with self.subTest(source_url=source_url):
                with self.assertRaises(ValueError):
                    normalize_public_web_url(source_url)

    @patch("app.socket.getaddrinfo")
    def test_web_source_url_validation_accepts_public_https(self, mock_getaddrinfo):
        mock_getaddrinfo.return_value = [
            (2, 1, 6, "", ("93.184.216.34", 443)),
        ]

        normalized = normalize_public_web_url("https://Example.com/research#results")

        self.assertEqual(normalized, "https://example.com/research")

    @patch("app.socket.getaddrinfo")
    @patch("app.requests.Session")
    def test_web_source_fetch_extracts_visible_html_only(self, mock_session_class, mock_getaddrinfo):
        mock_getaddrinfo.return_value = [(2, 1, 6, "", ("93.184.216.34", 443))]
        response = Mock()
        response.status_code = 200
        response.headers = {"Content-Type": "text/html; charset=utf-8"}
        response.iter_content.return_value = [
            b"<html><head><title>Useful study</title><script>ignore me</script></head>",
            b"<body><h1>Finding</h1><p>Visible evidence.</p></body></html>",
        ]
        session = mock_session_class.return_value
        session.get.return_value = response

        source = fetch_web_source("https://example.com/research")

        self.assertEqual(source["title"], "Useful study")
        self.assertIn("Visible evidence.", source["text"])
        self.assertNotIn("ignore me", source["text"])
        self.assertFalse(session.trust_env)
        session.get.assert_called_once()
        response.close.assert_called_once()
        session.close.assert_called_once()

    @patch("app.socket.getaddrinfo")
    @patch("app.requests.Session")
    def test_web_source_fetch_blocks_redirect_to_private_network(self, mock_session_class, mock_getaddrinfo):
        mock_getaddrinfo.return_value = [(2, 1, 6, "", ("93.184.216.34", 443))]
        response = Mock()
        response.status_code = 302
        response.headers = {"Location": "http://127.0.0.1/admin"}
        session = mock_session_class.return_value
        session.get.return_value = response

        with self.assertRaisesRegex(ValueError, "Local and private"):
            fetch_web_source("https://example.com/redirect")

        session.get.assert_called_once()
        response.close.assert_called_once()
        session.close.assert_called_once()

    def test_settings_api_saves_provider_choice(self):
        response = self.client.post(
            "/api/settings",
            json={
                "MODEL_PROVIDER": "demo",
                "OLLAMA_BASE_URL": "http://127.0.0.1:11434",
                "OLLAMA_MODEL": "qwen3.5:9b",
                "OPENAI_API_KEY": "",
                "OPENAI_MODEL": "gpt-5-mini",
                "ANTHROPIC_API_KEY": "",
                "CLAUDE_MODEL": "claude-sonnet-4-5",
            },
        )
        self.assertEqual(response.status_code, 200)
        payload = response.get_json()
        self.assertTrue(payload["ok"])
        self.assertEqual(payload["model_provider"], "demo")
        self.assertIn(b"Model settings", self.client.get("/").data)

    def test_local_model_selector_offers_qwen_and_phi_in_both_workspaces(self):
        for path in ("/", "/prompt-library"):
            with self.subTest(path=path):
                response = self.client.get(path)
                self.assertEqual(response.status_code, 200)
                self.assertIn(b'<select name="OLLAMA_MODEL" data-ollama-model>', response.data)
                self.assertIn(b'value="qwen3.5:4b"', response.data)
                self.assertIn(b'value="phi4-mini"', response.data)
                self.assertIn(b"best scientific critique", response.data)
                self.assertIn(b"faster responses", response.data)

    def test_settings_api_can_switch_to_phi4_mini(self):
        response = self.client.post(
            "/api/settings",
            json={
                "MODEL_PROVIDER": "ollama",
                "OLLAMA_BASE_URL": "http://127.0.0.1:11434",
                "OLLAMA_MODEL": "phi4-mini",
                "OPENAI_API_KEY": "",
                "OPENAI_MODEL": "gpt-5-mini",
                "ANTHROPIC_API_KEY": "",
                "CLAUDE_MODEL": "claude-sonnet-4-5",
            },
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.get_json()["settings"]["OLLAMA_MODEL"], "phi4-mini")
        self.assertEqual(response.get_json()["provider_label"], "phi4-mini")
        self.assertEqual(self.client.get("/api/settings").get_json()["OLLAMA_MODEL"], "phi4-mini")
        home = self.client.get("/")
        self.assertIn(b'<option value="phi4-mini" selected>', home.data)
        self.assertIn(b'<span data-provider-status-text>phi4-mini</span>', home.data)
        prompt_library = self.client.get("/prompt-library")
        self.assertIn(b'<span data-provider-status-text>phi4-mini</span>', prompt_library.data)

    def test_installer_selects_qwen_then_optionally_adds_phi(self):
        installer = (Path("installer") / "Promptly-Setup.ps1").read_text(encoding="utf-8")

        self.assertIn('Read-Host "Enter 1, 2, or 3"', installer)
        self.assertIn('$models = @($model)', installer)
        self.assertIn('Also install Phi-4 Mini so you can switch models in the website?', installer)
        self.assertIn('$models += "phi4-mini"', installer)
        self.assertNotIn('Enter 1, 2, 3, 4, or 5', installer)
        self.assertIn("foreach ($downloadModel in $models)", installer)
        self.assertIn("OLLAMA_MODEL=$model", installer)
        self.assertIn("Installed model(s): $($models -join ', ')", installer)
        self.assertIn("promptly-icon.ico", installer)
        self.assertIn("$shortcut.IconLocation", installer)
        self.assertIn("Stop-PromptlyProcesses $InstallDirectory", installer)
        self.assertIn("Remove-ExistingPrivateEnvironment $InstallDirectory", installer)
        self.assertIn("Stopping an existing Promptly installation", installer)
        self.assertIn("could not replace its existing private Python environment", installer)

    def test_promptly_icon_assets_are_valid_and_used_by_both_pages(self):
        png_path = Path("static") / "promptly-icon.png"
        ico_path = Path("static") / "promptly-icon.ico"
        self.assertTrue(png_path.is_file())
        self.assertTrue(ico_path.is_file())

        with Image.open(png_path) as png:
            self.assertEqual(png.size, (1024, 1024))
            self.assertEqual(png.mode, "RGBA")
            self.assertEqual(png.getpixel((0, 0))[3], 0)

        with Image.open(ico_path) as icon:
            self.assertEqual(icon.format, "ICO")
            self.assertIn((16, 16), icon.ico.sizes())
            self.assertIn((256, 256), icon.ico.sizes())

        for path in ("/", "/prompt-library"):
            with self.subTest(path=path):
                response = self.client.get(path)
                self.assertEqual(response.status_code, 200)
                self.assertIn(b'rel="icon"', response.data)
                self.assertIn(b'/static/promptly-icon.png', response.data)

    def test_rejects_wrong_file_type(self):
        response = self.client.post(
            "/feedback",
            data={"content_type": "papers-proposals", "file": (BytesIO(b"not a paper"), "notes.pptx")},
            content_type="multipart/form-data",
        )
        self.assertEqual(response.status_code, 400)
        self.assertIn(b"not supported", response.data)

    def test_damaged_feedback_files_return_clear_validation_error(self):
        for content_type, filename in (
            ("papers-proposals", "broken.pdf"),
            ("papers-proposals", "broken.docx"),
            ("talks-slides", "broken.pptx"),
        ):
            with self.subTest(filename=filename):
                response = self.client.post(
                    "/feedback",
                    data={
                        "content_type": content_type,
                        "mentor_id": "dr-nanshu-lu",
                        "file": (BytesIO(b"not a real file"), filename),
                    },
                    content_type="multipart/form-data",
                )
                self.assertEqual(response.status_code, 400)
                self.assertIn(b"could not be read", response.data)

    def test_damaged_reference_file_returns_clear_validation_error(self):
        response = self.client.post(
            "/generate-prompts",
            data={
                "prompt_mentor_name": "Damaged Reference",
                "research_files": (BytesIO(b"not a real PDF"), "broken.pdf"),
            },
            content_type="multipart/form-data",
        )
        self.assertEqual(response.status_code, 400)
        self.assertIn(b"could not be read", response.data)

    def test_rejects_unknown_mentor(self):
        response = self.client.post(
            "/feedback",
            data={
                "content_type": "papers-proposals",
                "mentor_id": "unknown-mentor",
                "file": (BytesIO(b"Example content"), "example.txt"),
            },
            content_type="multipart/form-data",
        )
        self.assertEqual(response.status_code, 400)
        self.assertIn(b"available mentor", response.data)

    def test_home_page_lists_library_mentors_for_feedback(self):
        create = self.client.post(
            "/generate-prompts",
            data={
                "prompt_mentor_name": "Custom Lab",
                "slide_files": (
                    BytesIO(b"Remove duplicate panels and clarify the takeaway."),
                    "slides.txt",
                ),
            },
            content_type="multipart/form-data",
        )
        self.assertEqual(create.status_code, 200)

        response = self.client.get("/?prompt_mentor=custom-lab&mentor=custom-lab&type=talks-slides")
        self.assertEqual(response.status_code, 200)
        self.assertIn(b"Choose a mentor for feedback", response.data)
        self.assertIn(b"Custom Lab", response.data)
        self.assertIn(b'name="mentor_id" value="custom-lab"', response.data)

    def test_feedback_uses_generated_library_prompt_for_matching_category(self):
        create = self.client.post(
            "/generate-prompts",
            data={
                "prompt_mentor_name": "Style Mentor",
                "slide_files": (
                    BytesIO(b"Remove duplicate panels and clarify the takeaway."),
                    "slides.txt",
                ),
            },
            content_type="multipart/form-data",
        )
        self.assertEqual(create.status_code, 200)

        # Category without its own generated prompt should fall back to the
        # library's available style prompt instead of blocking feedback.
        fallback = self.client.post(
            "/feedback",
            data={
                "content_type": "research-ideas",
                "mentor_id": "style-mentor",
                "file": (BytesIO(b"A research idea."), "idea.txt"),
            },
            content_type="multipart/form-data",
        )
        self.assertEqual(fallback.status_code, 200)
        self.assertIn(b"uploaded and read successfully", fallback.data)

        os.environ["MODEL_PROVIDER"] = "ollama"
        with patch("app.call_ollama", return_value="Library mentor feedback") as mock_ollama:
            response = self.client.post(
                "/feedback",
                data={
                    "content_type": "talks-slides",
                    "mentor_id": "style-mentor",
                    "file": (BytesIO(b"Slide one: main claim"), "deck.txt"),
                },
                content_type="multipart/form-data",
            )
        self.assertEqual(response.status_code, 200)
        self.assertIn(b"Library mentor feedback", response.data)
        self.assertEqual(mock_ollama.call_count, 1)
        system_prompt = mock_ollama.call_args.args[0]
        self.assertIn("Style Mentor", system_prompt)
        self.assertIn("Use the professor's style", system_prompt)

        with patch("app.call_ollama", return_value="Fallback style feedback") as mock_fallback:
            mismatched = self.client.post(
                "/feedback",
                data={
                    "content_type": "research-ideas",
                    "mentor_id": "style-mentor",
                    "file": (BytesIO(b"A research idea needing review."), "idea2.txt"),
                },
                content_type="multipart/form-data",
            )
        self.assertEqual(mismatched.status_code, 200)
        self.assertIn(b"Fallback style feedback", mismatched.data)
        fallback_system = mock_fallback.call_args.args[0]
        self.assertIn("No Research Ideas / Meeting Minutes prompt", fallback_system)
        self.assertIn("Talks / Presentations / Slides", fallback_system)

    def test_feedback_prompt_identifies_mentor(self):
        prompt = build_feedback_prompt(
            "papers-proposals",
            "example.txt",
            "Example content",
            "clarity",
            "dr-nanshu-lu",
        )
        self.assertIn("Dr. Nanshu Lu", prompt)
        self.assertIn("specialty, thinking process, and feedback style", prompt)
        self.assertIn("strength of the argument", prompt)
        self.assertIn("## Critical questions", prompt)
        self.assertIn("do not number questions as new sections", prompt)

    def test_feedback_markdown_repairs_math_sections_and_nested_questions(self):
        feedback = r"""6. Critical questions

7. What are the sample sizes ($n$)?

8. Report error bars ($\pm$ SD).

9. Is the result significant ($p > 0.05$)?

10. Is the fabrication complexity justified?

11. Prioritized revisions

12. Add sample sizes.

13. Add error bars."""

        rendered = render_markdown_html(feedback)

        self.assertIn("<h2>Critical questions</h2>", rendered)
        self.assertIn("<li>What are the sample sizes (n)?</li>", rendered)
        self.assertIn("<li>Report error bars (± SD).</li>", rendered)
        self.assertIn("<li>Is the result significant (p &gt; 0.05)?</li>", rendered)
        self.assertIn("<h2>Prioritized revisions</h2>", rendered)
        self.assertIn("<ol>", rendered)
        self.assertNotIn('<ol start="6">', rendered)
        self.assertNotIn("$n$", rendered)
        self.assertNotIn(r"\pm", rendered)

    def test_feedback_markdown_handles_long_non_heading_line(self):
        long_line = "Unrecognized heading " + (" " * 100_000) + "!"

        rendered = render_markdown_html(long_line)

        self.assertIn("Unrecognized heading", rendered)
        self.assertIn("!", rendered)

    def test_feedback_markdown_renders_inline_and_display_latex_as_readable_text(self):
        feedback = r"""## Technical corrections needed

The normalized error is \[
E = \frac{\Delta L}{L_0} \times 100\%.
\]

Require \(x_i \leq \sqrt{\alpha^2 + \beta^2}\), and report
$$p_{max} \approx 3 \times 10^{-4}.$$
"""

        rendered = render_markdown_html(feedback)

        self.assertIn("E = (Δ L)/(L₀) × 100%.", rendered)
        self.assertIn("xᵢ ≤ √(α² + β²)", rendered)
        self.assertIn("p_(max) ≈ 3 × 10⁻⁴", rendered)
        for raw_latex in (r"\[", r"\]", "$$", r"\frac", r"\sqrt", r"\alpha"):
            self.assertNotIn(raw_latex, rendered)

    def test_feedback_markdown_does_not_treat_currency_as_math(self):
        rendered = render_markdown_html("The first option costs $5 and the second costs $10.")

        self.assertIn("$5 and the second costs $10", rendered)

    def test_mentor_style_prompts_are_available_for_all_three_modes(self):
        research_prompt = load_mentor_prompt("dr-nanshu-lu", "research-ideas")
        paper_prompt = load_mentor_prompt("dr-nanshu-lu", "papers-proposals")
        slides_prompt = load_mentor_prompt("dr-nanshu-lu", "talks-slides")
        self.assertIn("falsifiable", research_prompt)
        self.assertIn("central claim", paper_prompt)
        self.assertIn("central story", slides_prompt)
        self.assertNotEqual(research_prompt, paper_prompt)
        self.assertNotEqual(paper_prompt, slides_prompt)

    @patch("app.OpenAI")
    def test_openai_receives_mentor_profile_and_review(self, mock_openai):
        os.environ["MODEL_PROVIDER"] = "openai"
        os.environ["OPENAI_API_KEY"] = "test-key"
        mock_openai.return_value.responses.create.return_value = SimpleNamespace(
            output_text="Focused mentor feedback"
        )

        result = call_model(
            "Review category: Research ideas\nContent: A testable idea.",
            mentor_id="dr-nanshu-lu",
        )

        self.assertEqual(result, "Focused mentor feedback")
        request = mock_openai.return_value.responses.create.call_args.kwargs
        self.assertEqual(request["model"], "gpt-5-mini")
        self.assertIn("Dr. Nanshu Lu", request["instructions"])
        self.assertIn("rigorous senior research mentor", request["instructions"])
        self.assertIn("A testable idea", request["input"])
        self.assertFalse(request["store"])

    @patch("app.requests.post")
    def test_ollama_uses_bounded_non_thinking_request_locally(self, mock_post):
        os.environ["MODEL_PROVIDER"] = "ollama"
        os.environ["OLLAMA_MODEL"] = "qwen3.5:9b"
        mock_post.return_value.json.return_value = {
            "message": {"content": "Critical local feedback", "thinking": ""}
        }

        result = call_model(
            "Review category: Research ideas\nContent: A testable idea.",
            mentor_id="dr-nanshu-lu",
        )

        self.assertEqual(result, "Critical local feedback")
        request = mock_post.call_args
        self.assertEqual(request.args[0], "http://127.0.0.1:11434/api/chat")
        payload = request.kwargs["json"]
        self.assertEqual(payload["model"], "qwen3.5:9b")
        self.assertFalse(payload["think"])
        self.assertFalse(payload["stream"])
        self.assertEqual(payload["keep_alive"], "30m")
        self.assertEqual(payload["options"]["num_predict"], 1_800)
        self.assertEqual(payload["options"]["num_ctx"], 8_192)
        self.assertEqual(request.kwargs["timeout"], 300)
        self.assertIn("Dr. Nanshu Lu", payload["messages"][0]["content"])
        self.assertIn("A testable idea", payload["messages"][1]["content"])

    @patch("app.requests.post")
    def test_ollama_does_not_repeat_a_slow_request_when_output_is_empty(self, mock_post):
        os.environ["MODEL_PROVIDER"] = "ollama"
        os.environ["OLLAMA_MODEL"] = "qwen3.5:9b"
        mock_post.return_value.json.return_value = {
            "message": {"content": "", "thinking": ""}
        }

        with self.assertRaisesRegex(ValueError, "Qwen 4B"):
            call_model(
                "Review category: Research ideas\nContent: A testable idea.",
                mentor_id="dr-nanshu-lu",
            )

        self.assertEqual(mock_post.call_count, 1)
        self.assertFalse(mock_post.call_args.kwargs["json"]["think"])

    def test_long_review_text_is_split_without_losing_content(self):
        text = ("A research claim with supporting evidence. " * 3_000).strip()
        chunks = split_review_text(text, max_chars=10_000)
        self.assertGreater(len(chunks), 1)
        self.assertTrue(all(len(chunk) <= 10_000 for chunk in chunks))
        self.assertEqual(" ".join(chunks), text)

    @patch("app.requests.post")
    def test_ollama_reviews_large_files_in_chunks_then_synthesizes(self, mock_post):
        os.environ["MODEL_PROVIDER"] = "ollama"
        os.environ["OLLAMA_MODEL"] = "qwen3.5:9b"
        responses = []
        for content in ("First section review", "Second section review", "Final synthesis"):
            response = Mock()
            response.json.return_value = {"message": {"content": content, "thinking": ""}}
            responses.append(response)
        mock_post.side_effect = responses
        long_prompt = "A" * 75_000

        result = call_model(long_prompt, mentor_id="dr-nanshu-lu")

        self.assertEqual(result, "Final synthesis")
        self.assertEqual(mock_post.call_count, 3)
        first_prompt = mock_post.call_args_list[0].kwargs["json"]["messages"][1]["content"]
        final_prompt = mock_post.call_args_list[2].kwargs["json"]["messages"][1]["content"]
        self.assertIn("section 1 of 2", first_prompt)
        self.assertIn("First section review", final_prompt)
        self.assertIn("Second section review", final_prompt)
        self.assertNotIn("A" * 1_000, final_prompt)
        partial_payload = mock_post.call_args_list[0].kwargs["json"]
        final_payload = mock_post.call_args_list[2].kwargs["json"]
        self.assertFalse(partial_payload["think"])
        self.assertEqual(partial_payload["options"]["num_predict"], 700)
        self.assertEqual(final_payload["options"]["num_predict"], 1_800)
        self.assertLessEqual(partial_payload["options"]["num_ctx"], 16_384)
        self.assertLessEqual(final_payload["options"]["num_ctx"], 16_384)

    @patch("app.requests.post")
    def test_ollama_caps_large_review_at_three_sections_plus_synthesis(self, mock_post):
        os.environ["MODEL_PROVIDER"] = "ollama"
        responses = []
        for content in ("Beginning review", "Middle review", "Ending review", "Final synthesis"):
            response = Mock()
            response.json.return_value = {"message": {"content": content, "thinking": ""}}
            responses.append(response)
        mock_post.side_effect = responses

        result = call_model("A" * 300_000, mentor_id="dr-nanshu-lu")

        self.assertEqual(result, "Final synthesis")
        self.assertEqual(mock_post.call_count, 4)
        synthesis = mock_post.call_args_list[-1].kwargs["json"]["messages"][1]["content"]
        self.assertIn("Beginning review", synthesis)
        self.assertIn("Middle review", synthesis)
        self.assertIn("Ending review", synthesis)

    @patch("app.Anthropic")
    def test_claude_receives_mentor_profile_and_review(self, mock_anthropic):
        os.environ["MODEL_PROVIDER"] = "claude"
        os.environ["ANTHROPIC_API_KEY"] = "test-key"
        mock_anthropic.return_value.messages.create.return_value = SimpleNamespace(
            content=[SimpleNamespace(type="text", text="Focused Claude mentor feedback")]
        )

        result = call_model(
            "Review category: Research ideas\nContent: A testable idea.",
            mentor_id="dr-nanshu-lu",
        )

        self.assertEqual(result, "Focused Claude mentor feedback")
        request = mock_anthropic.return_value.messages.create.call_args.kwargs
        self.assertEqual(request["model"], "claude-sonnet-4-5")
        self.assertEqual(request["max_tokens"], 2_000)
        self.assertIn("Dr. Nanshu Lu", request["system"])
        self.assertIn("rigorous senior research mentor", request["system"])
        self.assertEqual(request["messages"][0]["role"], "user")
        self.assertIn("A testable idea", request["messages"][0]["content"])

    def test_word_document_upload_is_read(self):
        file_data = BytesIO()
        document = Document()
        document.add_paragraph("A proposal that needs feedback.")
        document.save(file_data)
        file_data.seek(0)
        response = self.client.post(
            "/feedback",
            data={"content_type": "papers-proposals", "file": (file_data, "proposal.docx")},
            content_type="multipart/form-data",
        )
        self.assertEqual(response.status_code, 200)
        self.assertIn(b"uploaded and read successfully", response.data)

    def test_powerpoint_upload_is_read(self):
        file_data = BytesIO()
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Project update"
        slide.placeholders[1].text = "The project is on schedule."
        presentation.save(file_data)
        file_data.seek(0)
        response = self.client.post(
            "/feedback",
            data={"content_type": "talks-slides", "file": (file_data, "update.pptx")},
            content_type="multipart/form-data",
        )
        self.assertEqual(response.status_code, 200)
        self.assertIn(b"uploaded and read successfully", response.data)

    def test_json_api_validates_empty_prompt(self):
        response = self.client.post("/api/generate", json={"prompt": ""})
        self.assertEqual(response.status_code, 400)


if __name__ == "__main__":
    unittest.main()
